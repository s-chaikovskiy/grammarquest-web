# -*- coding: utf-8 -*-
"""
Движок вёрстки премиальных презентаций.
Ключевая идея: весь текст измеряется по реальным метрикам TTF через PIL,
строки разбиваются заранее, высота блока вычисляется точно.
Наезд заголовка на подзаголовок физически невозможен.
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import ImageFont
import copy

FONT_DIR = os.path.expanduser("~/Library/Fonts")
EMU_IN = 914400.0

# ─────────────────────────────────────── шрифты
# Гарнитура та же, что в самом приложении: Golos Text. Выбрана потому, что
# полностью покрывает казахские буквы ә ғ қ ң ө ұ ү һ і — у большинства
# популярных шрифтов их нет, и на слайде они подменились бы чужим начертанием.
DISPLAY_XB = "GQ Text Bold"
DISPLAY_B = "GQ Text Bold"
DISPLAY_M = "GQ Text SemiBold"
TEXT = "GQ Text"
TEXT_MD = "GQ Text Medium"
TEXT_SB = "GQ Text SemiBold"
TEXT_IT = "GQ Text"
MONO = "IBM Plex Mono"
MONO_MD = "IBM Plex Mono Medium"
MONO_SB = "IBM Plex Mono SemiBold"

_FILE = {
    DISPLAY_XB: "GQText-Bold.ttf",
    DISPLAY_B: "GQText-Bold.ttf",
    DISPLAY_M: "GQText-SmBd.ttf",
    TEXT: "GQText-Rg.ttf",
    TEXT_MD: "GQText-Md.ttf",
    TEXT_SB: "GQText-SmBd.ttf",
    TEXT_IT: "GQText-Rg.ttf",
    MONO: "IBMPlexMono-Regular.ttf",
    MONO_MD: "IBMPlexMono-Medium.ttf",
    MONO_SB: "IBMPlexMono-SemiBold.ttf",
}

# ─────────────────────────────────────── палитра
# Те же цвета, что в приложении: слайды и продукт должны выглядеть
# как одна работа, а не как доклад отдельно и программа отдельно.
PAPER = RGBColor(0xF6, 0xF7, 0xF9)     # фон светлого слайда
PAPER_2 = RGBColor(0xEC, 0xF1, 0xF5)   # заливка блоков
PAPER_3 = RGBColor(0xDD, 0xE4, 0xEA)   # плотная заливка
INK = RGBColor(0x1B, 0x22, 0x2B)       # основной текст
INK_2 = RGBColor(0x11, 0x16, 0x1E)     # фон тёмного слайда
SOFT = RGBColor(0x53, 0x5C, 0x66)      # вторичный текст на светлом
SOFT_D = RGBColor(0xA6, 0xAD, 0xB6)    # вторичный текст на тёмном
HAIR = RGBColor(0xCC, 0xD4, 0xDC)      # волосяная линия на светлом
HAIR_D = RGBColor(0x3B, 0x43, 0x4D)    # волосяная линия на тёмном
SIGNAL = RGBColor(0x00, 0x64, 0xB9)    # акцент приложения
SIGNAL_L = RGBColor(0x59, 0xAA, 0xF8)  # акцент на тёмном фоне
PETROL = RGBColor(0x00, 0x79, 0x3D)    # «верно» — зелёный
PETROL_L = RGBColor(0x48, 0xCD, 0x8C)
OCHRE = RGBColor(0x7F, 0x54, 0x00)     # «серия, награда» — золото
OCHRE_L = RGBColor(0xEE, 0xBB, 0x58)
ERROR = RGBColor(0xBA, 0x1D, 0x27)     # «ошибка» — красный
ERROR_L = RGBColor(0xF6, 0x6D, 0x67)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER_ON_D = RGBColor(0xEE, 0xEB, 0xE5)

# ─────────────────────────────────────── сетка (дюймы)
W, H = 13.333, 7.5
ML, MR = 0.82, 0.82
MT, MB = 0.52, 0.46
CW = W - ML - MR            # 11.693
COLS = 12
GUT = 0.22
COLW = (CW - GUT * (COLS - 1)) / COLS


def col(n_from, n_span):
    """левый край и ширина для колонок сетки"""
    x = ML + n_from * (COLW + GUT)
    w = n_span * COLW + (n_span - 1) * GUT
    return x, w


# ─────────────────────────────────────── измерение
_cache = {}


def _font(name, size_pt):
    key = (name, round(size_pt * 4))
    if key not in _cache:
        path = os.path.join(FONT_DIR, _FILE[name])
        _cache[key] = ImageFont.truetype(path, int(round(size_pt * 4)))
    return _cache[key]


def text_w(s, name, size_pt, tracking=0.0):
    """ширина строки в дюймах; tracking в em"""
    if not s:
        return 0.0
    f = _font(name, size_pt)
    w = f.getlength(s) / 4.0
    w += tracking * size_pt * max(len(s) - 1, 0)
    return w / 72.0


def wrap(s, name, size_pt, max_in, tracking=0.0):
    """разбивка на строки по реальной ширине глифов; уважает явные \n"""
    out = []
    for para in s.split("\n"):
        words = para.split(" ")
        line = ""
        for wd in words:
            trial = wd if not line else line + " " + wd
            if text_w(trial, name, size_pt, tracking) <= max_in or not line:
                line = trial
            else:
                out.append(line)
                line = wd
        out.append(line)
    return out


def fit(s, name, max_in, max_lines, start_pt, min_pt, step=0.5, tracking=0.0):
    """
    Подбирает кегль так, чтобы текст уложился в max_lines строк И по ширине.
    Проверка ширины обязательна: wrap не разрывает одиночное длинное слово,
    поэтому без неё длинное слово «помещается» на любом кегле и вылезает за блок.
    """
    size = start_pt
    while size > min_pt:
        lines = wrap(s, name, size, max_in, tracking)
        widest = max((text_w(l, name, size, tracking) for l in lines), default=0.0)
        if len(lines) <= max_lines and widest <= max_in:
            return size, lines
        size -= step
    return min_pt, wrap(s, name, min_pt, max_in, tracking)


# ─────────────────────────────────────── низкоуровневая укладка
class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Emu(int(W * EMU_IN))
        self.prs.slide_height = Emu(int(H * EMU_IN))
        self.blank = self.prs.slide_layouts[6]
        self.slides = []
        self.total = 0

    def add(self, dark=False):
        s = self.prs.slides.add_slide(self.blank)
        bg = PAPER if not dark else INK
        r = rect(s, 0, 0, W, H, bg)
        self.slides.append(s)
        return s

    def save(self, path):
        # проставляем «NN / total»
        self.prs.save(path)



def _strip_style(shape):
    """убирает <p:style> — иначе LibreOffice подмешивает тень и обводку темы"""
    el = shape._element
    st = el.find(qn("p:style"))
    if st is not None:
        el.remove(st)


def _noline(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False


def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75, radius=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(x * EMU_IN)), Emu(int(y * EMU_IN)),
                                 Emu(int(w * EMU_IN)), Emu(int(h * EMU_IN)))
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    _strip_style(shp)
    shp.text_frame.word_wrap = False
    return shp


def hline(slide, x, y, w, color=HAIR, lw=0.75):
    return rect(slide, x, y, w, 0.008, fill=color)


def vline(slide, x, y, h, color=HAIR, lw=0.75):
    return rect(slide, x, y, 0.008, h, fill=color)


def dot(slide, x, y, d, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(x * EMU_IN)), Emu(int(y * EMU_IN)),
                                 Emu(int(d * EMU_IN)), Emu(int(d * EMU_IN)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    _noline(shp)
    _strip_style(shp)
    shp.text_frame.word_wrap = False
    return shp


def _set_spacing(p, pts):
    pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
    for tag in ("a:lnSpc",):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    from lxml import etree
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    pts_el = etree.SubElement(lnSpc, qn("a:spcPts"))
    pts_el.set("val", str(int(round(pts * 100))))
    pPr.insert(0, lnSpc)


def _set_tracking(run, em, size_pt):
    if not em:
        return
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(round(em * size_pt * 100))))


def block(slide, lines, x, y, w, name, size_pt, color, leading=None,
          tracking=0.0, align=PP_ALIGN.LEFT, caps=False, space_after=0.0):
    """
    Кладёт заранее разбитые строки. Возвращает нижнюю границу (дюймы).
    leading — межстрочный в pt (по умолчанию 1.22 кегля).
    """
    if isinstance(lines, str):
        lines = [lines]
    lines = [l for l in lines]
    if not lines:
        return y
    lead = leading if leading else size_pt * 1.22
    h = (lead * len(lines) + lead * 0.42) / 72.0
    tb = slide.shapes.add_textbox(Emu(int(x * EMU_IN)), Emu(int((y - lead * 0.16 / 72.0) * EMU_IN)),
                                  Emu(int(w * EMU_IN)), Emu(int(h * EMU_IN)))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        _set_spacing(p, lead)
        r = p.add_run()
        r.text = ln.upper() if caps else ln
        r.font.name = name
        r.font.size = Pt(size_pt)
        r.font.color.rgb = color
        _set_tracking(r, tracking, size_pt)
        _latin(r, name)
    return y + (lead * len(lines)) / 72.0 + space_after


def _latin(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs", "a:ea"):
        el = rPr.find(qn(tag))
        if el is None:
            from lxml import etree
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


def rich(slide, segments, x, y, w, size_pt, leading=None, align=PP_ALIGN.LEFT):
    """одна строка из разноформатных кусков: [(text, font, color, tracking), ...]"""
    lead = leading if leading else size_pt * 1.22
    h = (lead * 1.5) / 72.0
    tb = slide.shapes.add_textbox(Emu(int(x * EMU_IN)), Emu(int((y - lead * 0.16 / 72.0) * EMU_IN)),
                                  Emu(int(w * EMU_IN)), Emu(int(h * EMU_IN)))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    _set_spacing(p, lead)
    for seg in segments:
        txt, fname, colr = seg[0], seg[1], seg[2]
        trk = seg[3] if len(seg) > 3 else 0.0
        sz = seg[4] if len(seg) > 4 else size_pt
        r = p.add_run()
        r.text = txt
        r.font.name = fname
        r.font.size = Pt(sz)
        r.font.color.rgb = colr
        _set_tracking(r, trk, sz)
        _latin(r, fname)
    return y + lead / 72.0


# ─────────────────────────────────────── компоненты слайда
def chrome(slide, deck, eyebrow, idx, dark=False, section=None):
    """верхняя служебная линейка + номер слайда. Возвращает y контента."""
    c_hair = HAIR_D if dark else HAIR
    c_mono = SOFT_D if dark else SOFT
    y = MT
    hline(slide, ML, y, CW, c_hair)
    ty = y + 0.115
    if eyebrow:
        block(slide, [eyebrow], ML, ty, CW * 0.7, MONO_SB, 8.5, SIGNAL_L if dark else SIGNAL,
              leading=10.5, tracking=0.16, caps=True)
    num = f"{idx:02d} / {deck.total:02d}"
    nw = text_w(num, MONO_MD, 8.5, 0.10) + 0.02
    block(slide, [num], ML + CW - nw, ty, nw + 0.1, MONO_MD, 8.5, c_mono,
          leading=10.5, tracking=0.10)
    return y + 0.44


def title(slide, text, y, dark=False, w=None, max_lines=2, start=40, minimum=25, color=None):
    w = w if w else CW
    size, lines = fit(text, DISPLAY_XB, w, max_lines, start, minimum, 0.5, -0.022)
    lead = size * 1.07
    c = color if color else (PAPER_ON_D if dark else INK)
    return block(slide, lines, ML, y, w, DISPLAY_XB, size, c, leading=lead, tracking=-0.022)


def lede(slide, text, y, w=None, dark=False, size=13.0, x=None, color=None, maxw=None):
    x = x if x is not None else ML
    w = w if w else CW * 0.78
    lines = wrap(text, TEXT, size, w)
    c = color if color else (SOFT_D if dark else SOFT)
    return block(slide, lines, x, y, w, TEXT, size, c, leading=size * 1.42)


def footnote(slide, text, dark=False, y=None, color=None):
    y = y if y else H - MB - 0.14
    lines = wrap(text, TEXT, 8.5, CW)
    c = color if color else (SOFT_D if dark else SOFT)
    yy = y - (len(lines) - 1) * 0.145
    return block(slide, lines, ML, yy, CW, TEXT, 8.5, c, leading=11.6)


# ─────────────────────────────────────── компоненты
OVERFLOW = []


def kpi_row(slide, items, y, dark=False, size=42, w_total=None, x0=None, gap=None,
            accent_first=False, label_size=8.5, sub_size=9.5):
    """items: [(число, подпись) | (число, подпись, пояснение)]"""
    x0 = x0 if x0 is not None else ML
    w_total = w_total if w_total else CW
    n = len(items)
    gap = gap if gap else 0.3
    cw = (w_total - gap * (n - 1)) / n
    c_num = PAPER_ON_D if dark else INK
    c_lbl = SOFT_D if dark else SOFT
    bottom = y
    for i, it in enumerate(items):
        x = x0 + i * (cw + gap)
        num, lbl = it[0], it[1]
        sub = it[2] if len(it) > 2 else None
        cn = SIGNAL_L if (dark and accent_first and i == 0) else (SIGNAL if (accent_first and i == 0) else c_num)
        sz = size
        while text_w(num, DISPLAY_XB, sz, -0.03) > cw and sz > 16:
            sz -= 1
        b = block(slide, [num], x, y, cw, DISPLAY_XB, sz, cn, leading=sz * 0.98, tracking=-0.03)
        b = block(slide, wrap(lbl, MONO_MD, label_size, cw, 0.08), x, b + 0.10, cw,
                  MONO_MD, label_size, c_lbl, leading=label_size * 1.5, tracking=0.08, caps=True)
        if sub:
            b = block(slide, wrap(sub, TEXT, sub_size, cw), x, b + 0.10, cw, TEXT, sub_size,
                      c_lbl, leading=sub_size * 1.42)
        bottom = max(bottom, b)
    return bottom


def bullets(slide, items, x, y, w, dark=False, size=11, gap=0.15, marker="—",
            color=None, mcolor=None, lead_mult=1.40, bold_head=False):
    """items: строка или (жирный зачин, продолжение)"""
    c = color if color else (PAPER_ON_D if dark else INK)
    mc = mcolor if mcolor else (SIGNAL_L if dark else SIGNAL)
    ind = text_w(marker, MONO_MD, size, 0) + 0.14
    yy = y
    for it in items:
        head, rest = (it, None) if isinstance(it, str) else it
        block(slide, [marker], x, yy, ind, MONO_MD, size, mc, leading=size * lead_mult)
        if rest is None and not bold_head:
            lines = wrap(head, TEXT, size, w - ind)
            yy = block(slide, lines, x + ind, yy, w - ind, TEXT, size, c, leading=size * lead_mult)
        else:
            segs = [(head, TEXT_SB, c)] + ([(rest, TEXT, c)] if rest else [])
            ml = mixed_wrap(segs, w - ind, size)
            yy = mixed_block(slide, ml, x + ind, yy, w - ind, size, leading=size * lead_mult)
        yy += gap
    return yy - gap



def mixed_wrap(segs, max_in, size_pt):
    """
    Перенос строки для текста из кусков РАЗНЫХ начертаний.
    Каждое слово меряется своим шрифтом — иначе жирный зачин съезжает.
    segs: [(текст, шрифт, цвет)]. Возврат: [[(текст, шрифт, цвет), ...], ...]
    """
    toks = []
    for txt, fn, cl in segs:
        parts = txt.split(" ")
        for i, wd in enumerate(parts):
            if wd == "" and i not in (0, len(parts) - 1):
                continue
            toks.append([wd, fn, cl, i < len(parts) - 1])
    lines, cur, curw = [], [], 0.0
    space_w = {}
    for wd, fn, cl, had_space in toks:
        if fn not in space_w:
            space_w[fn] = text_w(" ", fn, size_pt)
        ww = text_w(wd, fn, size_pt)
        need = ww + (space_w[fn] if cur else 0)
        if curw + need > max_in and cur:
            lines.append(cur)
            cur, curw = [[wd, fn, cl]], ww
        else:
            if cur and cur[-1][1] == fn and cur[-1][2] == cl:
                cur[-1][0] += (" " if curw else "") + wd
            else:
                cur.append([(" " if cur else "") + wd, fn, cl])
            curw += need
    if cur:
        lines.append(cur)
    return lines


def mixed_block(slide, lines, x, y, w, size_pt, leading=None, align=PP_ALIGN.LEFT):
    """кладёт результат mixed_wrap; возвращает нижнюю границу"""
    lead = leading if leading else size_pt * 1.4
    if not lines:
        return y
    tb = slide.shapes.add_textbox(Emu(int(x * EMU_IN)), Emu(int((y - lead * 0.16 / 72.0) * EMU_IN)),
                                  Emu(int(w * EMU_IN)), Emu(int(((len(lines) + 0.5) * lead) / 72.0 * EMU_IN)))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        _set_spacing(p, lead)
        for txt, fn, cl in ln:
            r = p.add_run(); r.text = txt
            r.font.name = fn; r.font.size = Pt(size_pt); r.font.color.rgb = cl
            _latin(r, fn)
    return y + len(lines) * lead / 72.0


def _rich_wrapped(slide, head, rest, x, y, w, size, color, lead_mult):
    """первый кусок полужирный, дальше обычный, с корректным переносом"""
    full = head + ((" " + rest) if rest else "")
    lines = wrap(full, TEXT, size, w)
    lead = size * lead_mult
    tb = slide.shapes.add_textbox(Emu(int(x * EMU_IN)), Emu(int((y - lead * 0.16 / 72.0) * EMU_IN)),
                                  Emu(int(w * EMU_IN)), Emu(int(((len(lines) + 0.5) * lead) / 72.0 * EMU_IN)))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    consumed = 0
    hl = len(head)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_spacing(p, lead)
        start, end = consumed, consumed + len(ln)
        if start < hl:
            cut_at = min(end, hl) - start
            r = p.add_run(); r.text = ln[:cut_at]
            r.font.name = TEXT_SB; r.font.size = Pt(size); r.font.color.rgb = color; _latin(r, TEXT_SB)
            if cut_at < len(ln):
                r2 = p.add_run(); r2.text = ln[cut_at:]
                r2.font.name = TEXT; r2.font.size = Pt(size); r2.font.color.rgb = color; _latin(r2, TEXT)
        else:
            r = p.add_run(); r.text = ln
            r.font.name = TEXT; r.font.size = Pt(size); r.font.color.rgb = color; _latin(r, TEXT)
        consumed = end + 1
    return y + len(lines) * lead / 72.0


def panel(slide, x, y, w, h, dark=False, fill=None, border=True, accent=None, accent_w=0.035):
    """плоский блок с волосяной рамкой; accent — цветная полоса слева"""
    f = fill if fill is not None else (INK_2 if dark else PAPER_2)
    b = (HAIR_D if dark else HAIR) if border else None
    shp = rect(slide, x, y, w, h, fill=f, line=b)
    if accent is not None:
        rect(slide, x, y, accent_w, h, fill=accent)
    return shp


def card(slide, x, y, w, h, tag, head, items, dark=False, fill=None, accent=None,
         head_size=15, item_size=10.5, pad=0.28, tag_color=None):
    panel(slide, x, y, w, h, dark=dark, fill=fill, accent=accent)
    iw = w - pad * 2
    yy = y + pad
    if tag:
        tc = tag_color if tag_color else (SIGNAL_L if dark else SIGNAL)
        yy = block(slide, [tag], x + pad, yy, iw, MONO_SB, 8.5, tc, leading=10.5, tracking=0.16, caps=True)
        yy += 0.13
    if head:
        hs, hlines = fit(head, DISPLAY_B, iw, 2, head_size, 11, 0.5, -0.015)
        yy = block(slide, hlines, x + pad, yy, iw, DISPLAY_B, hs, PAPER_ON_D if dark else INK,
                   leading=hs * 1.14, tracking=-0.015)
        yy += 0.17
    if items:
        yy = bullets(slide, items, x + pad, yy, iw, dark=dark, size=item_size, gap=0.11)
    need = yy - y + pad
    if need > h + 0.005:
        OVERFLOW.append(f"card «{(head or tag or '?')[:38]}»: нужно {need:.2f}\" при заданных {h:.2f}\" (не хватает {need - h:.2f}\")")
    return y + h


def table(slide, cols, rows, x, y, w, dark=False, head_size=8.5, row_size=10.5,
          row_h=0.42, align=None, mono_cols=(), accent_col=None, head_fill=None):
    """cols: [(заголовок, доля ширины)]; rows: [[...]]"""
    total = sum(c[1] for c in cols)
    widths = [w * c[1] / total for c in cols]
    xs, acc = [], x
    for cw_ in widths:
        xs.append(acc); acc += cw_
    c_hair = HAIR_D if dark else HAIR
    c_ink = PAPER_ON_D if dark else INK
    c_soft = SOFT_D if dark else SOFT
    align = align or [PP_ALIGN.LEFT] * len(cols)
    # шапка
    hh = 0.34
    if head_fill is not None:
        rect(slide, x, y, w, hh, fill=head_fill)
    for i, (t, _) in enumerate(cols):
        pad = 0.14 if head_fill is not None else 0.0
        block(slide, [t], xs[i] + pad, y + 0.085, widths[i] - pad * 2, MONO_SB, head_size,
              c_soft if head_fill is None else (PAPER_ON_D if dark else PAPER),
              leading=head_size * 1.3, tracking=0.14, caps=True, align=align[i])
    yy = y + hh
    hline(slide, x, yy, w, c_ink if head_fill is None else c_hair)
    yy += 0.02
    for r in rows:
        rh = row_h
        cells = []
        for i, cell in enumerate(r):
            fname = MONO_MD if i in mono_cols else (TEXT_SB if i == 0 else TEXT)
            pad = 0.14 if head_fill is not None else 0.0
            lines = wrap(str(cell), fname, row_size, widths[i] - pad * 2)
            cells.append((fname, lines, pad))
            rh = max(rh, len(lines) * row_size * 1.34 / 72.0 + 0.22)
        for i, (fname, lines, pad) in enumerate(cells):
            cc = c_ink
            if accent_col is not None and i == accent_col:
                cc = SIGNAL_L if dark else SIGNAL
                fname = MONO_SB if i in mono_cols else TEXT_SB
            block(slide, lines, xs[i] + pad, yy + (rh - len(lines) * row_size * 1.34 / 72.0) / 2 + 0.02,
                  widths[i] - pad * 2, fname, row_size, cc, leading=row_size * 1.34, align=align[i])
        yy += rh
        hline(slide, x, yy, w, c_hair)
        yy += 0.02
    return yy


def timeline(slide, stages, x, y, w, dark=False, h=None):
    """stages: [(номер, название, срок, [пункты])]"""
    n = len(stages)
    gap = 0.24
    cw = (w - gap * (n - 1)) / n
    c_hair = HAIR_D if dark else HAIR
    c_soft = SOFT_D if dark else SOFT
    bottom = y
    hline(slide, x, y, w, c_hair)
    for i, (num, name, term, pts) in enumerate(stages):
        xx = x + i * (cw + gap)
        rect(slide, xx, y - 0.005, cw, 0.03, fill=SIGNAL if i == 0 else (HAIR_D if dark else PAPER_3))
        yy = y + 0.24
        yy = block(slide, [num], xx, yy, cw, MONO_SB, 9, SIGNAL if not dark else SIGNAL_L,
                   leading=11, tracking=0.14, caps=True)
        yy += 0.09
        hs, hlines = fit(name, DISPLAY_B, cw, 2, 16, 12, 0.5, -0.015)
        yy = block(slide, hlines, xx, yy, cw, DISPLAY_B, hs, PAPER_ON_D if dark else INK,
                   leading=hs * 1.12, tracking=-0.015)
        yy += 0.08
        yy = block(slide, [term], xx, yy, cw, MONO_MD, 9, c_soft, leading=11.5, tracking=0.06)
        yy += 0.17
        yy = bullets(slide, pts, xx, yy, cw, dark=dark, size=10, gap=0.10, marker="·")
        bottom = max(bottom, yy)
    return bottom


def band(slide, y, text_segments, dark=False, fill=None, h=None, pad=0.30, size=12.0, x=None, w=None):
    """широкая акцентная полоса с текстом; text_segments — как в rich(), но с переносом"""
    x = x if x is not None else ML
    w = w if w is not None else CW
    f = fill if fill is not None else (SIGNAL if not dark else SIGNAL)
    segs = [(t[0], t[1], t[2]) for t in text_segments]
    lines = mixed_wrap(segs, w - pad * 2, size)
    hh = h if h else len(lines) * size * 1.45 / 72.0 + pad * 1.75
    rect(slide, x, y, w, hh, fill=f)
    mixed_block(slide, lines, x + pad, y + pad * 0.86, w - pad * 2, size, leading=size * 1.45)
    return y + hh


def _seg_wrapped(slide, segs, x, y, w, size, default_color, lead_mult):
    """многострочная раскладка разноформатных сегментов"""
    flat = []
    for s in segs:
        txt, fname, colr = s[0], s[1], s[2]
        for ch in txt:
            flat.append((ch, fname, colr))
    plain = "".join(c[0] for c in flat)
    lines = wrap(plain, TEXT, size, w)
    lead = size * lead_mult
    tb = slide.shapes.add_textbox(Emu(int(x * EMU_IN)), Emu(int((y - lead * 0.16 / 72.0) * EMU_IN)),
                                  Emu(int(w * EMU_IN)), Emu(int(((len(lines) + 0.5) * lead) / 72.0 * EMU_IN)))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    pos = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_spacing(p, lead)
        chunk = flat[pos:pos + len(ln)]
        pos += len(ln) + 1
        cur_f, cur_c, buf = None, None, ""
        for ch, fn, cl in chunk:
            if fn != cur_f or cl != cur_c:
                if buf:
                    r = p.add_run(); r.text = buf; r.font.name = cur_f
                    r.font.size = Pt(size); r.font.color.rgb = cur_c; _latin(r, cur_f)
                cur_f, cur_c, buf = fn, cl, ch
            else:
                buf += ch
        if buf:
            r = p.add_run(); r.text = buf; r.font.name = cur_f
            r.font.size = Pt(size); r.font.color.rgb = cur_c; _latin(r, cur_f)
    return y + len(lines) * lead / 72.0


# ─────────────────────────────────────── диаграммы (нативные векторные)
def hbars(slide, items, x, y, w, dark=False, bar_h=0.28, gap=0.22, label_w=None,
          val_w=0.95, colors=None, max_val=None, size=10.5):
    """items: [(подпись, значение, строка_значения)] — горизонтальные полосы"""
    c_soft = SOFT_D if dark else SOFT
    c_ink = PAPER_ON_D if dark else INK
    label_w = label_w if label_w else w * 0.32
    plot_x = x + label_w + 0.18
    plot_w = w - label_w - 0.18 - val_w
    mx = max_val if max_val else max(i[1] for i in items)
    yy = y
    for i, it in enumerate(items):
        lbl, val, vs = it[0], it[1], it[2]
        cl = colors[i] if colors else (SIGNAL if i == 0 else (PAPER_3 if not dark else HAIR_D))
        lines = wrap(lbl, TEXT, size, label_w)
        block(slide, lines, x, yy + (bar_h - len(lines) * size * 1.3 / 72.0) / 2 + 0.015,
              label_w, TEXT, size, c_ink, leading=size * 1.3)
        bw = max(plot_w * val / mx, 0.02) if mx else 0.02
        rect(slide, plot_x, yy, plot_w, bar_h, fill=PAPER_2 if not dark else INK_2)
        rect(slide, plot_x, yy, bw, bar_h, fill=cl)
        block(slide, [vs], plot_x + plot_w + 0.12, yy + (bar_h - size * 1.3 / 72.0) / 2 + 0.015,
              val_w, MONO_SB, size, c_ink, leading=size * 1.3, align=PP_ALIGN.RIGHT)
        rh = max(bar_h, len(lines) * size * 1.3 / 72.0)
        yy += rh + gap
    return yy - gap


def stacked(slide, groups, x, y, w, h, dark=False, colors=None, legend=None,
            fmt=None, gap=0.5, show_total=True, axis_label=None,
            seg_labels=False, seg_min=0.19, seg_size=10.0, seg_color=None):
    """
    groups: [(название, [сегмент1, сегмент2, ...])] — вертикальные столбцы с накоплением
    colors: [цвет на каждый сегмент]
    seg_labels — подписывать значение внутри сегмента. Сегмент ниже seg_min
    остаётся без подписи: втиснутая в три пикселя цифра хуже, чем её отсутствие.
    """
    c_soft = SOFT_D if dark else SOFT
    c_ink = PAPER_ON_D if dark else INK
    n = len(groups)
    bw = (w - gap * (n - 1)) / n
    totals = [sum(g[1]) for g in groups]
    mx = max(totals) * 1.0
    base = y + h
    for i, (name, segs) in enumerate(groups):
        xx = x + i * (bw + gap)
        cy = base
        for j, v in enumerate(segs):
            sh = (v / mx) * h if mx else 0
            if sh <= 0:
                continue
            cl = colors[j] if colors else [SIGNAL, OCHRE, PETROL, PAPER_3][j % 4]
            rect(slide, xx, cy - sh, bw, sh, fill=cl)
            if seg_labels and sh >= seg_min:
                lv = fmt(v) if fmt else str(v)
                block(slide, [lv], xx, cy - sh + (sh - seg_size * 1.05 / 72.0) / 2 - 0.02, bw,
                      MONO_SB, seg_size, seg_color if seg_color is not None else WHITE,
                      leading=seg_size * 1.05, align=PP_ALIGN.CENTER)
            cy -= sh
        if show_total:
            tv = fmt(totals[i]) if fmt else str(totals[i])
            sz = 24
            while text_w(tv, DISPLAY_XB, sz, -0.03) > bw and sz > 12:
                sz -= 1
            block(slide, [tv], xx, cy - sz * 1.05 / 72.0 - 0.13, bw, DISPLAY_XB, sz, c_ink,
                  leading=sz * 1.05, tracking=-0.03, align=PP_ALIGN.CENTER)
        lines = wrap(name, MONO_MD, 9, bw, 0.06)
        block(slide, lines, xx, base + 0.16, bw, MONO_MD, 9, c_soft, leading=12,
              tracking=0.06, caps=True, align=PP_ALIGN.CENTER)
    hline(slide, x, base, w, HAIR_D if dark else HAIR)
    _lab_lines = max(len(wrap(g[0], MONO_MD, 9, bw, 0.06)) for g in groups)
    bottom = base + 0.16 + _lab_lines * 12 / 72.0 + 0.20
    if legend:
        lx = x
        for j, lab in enumerate(legend):
            cl = colors[j] if colors else [SIGNAL, OCHRE, PETROL, PAPER_3][j % 4]
            rect(slide, lx, bottom + 0.045, 0.115, 0.115, fill=cl)
            tw = text_w(lab, TEXT, 9.5) + 0.02
            block(slide, [lab], lx + 0.19, bottom, tw + 0.1, TEXT, 9.5, c_ink, leading=12.5)
            lx += 0.19 + tw + 0.34
        bottom += 0.30
    return bottom


def lines_chart(slide, series, x, y, w, h, dark=False, xlabels=None, ylab_fmt=None,
                ymax=None, size=9.5):
    """series: [(название, цвет, [значения])] — линейный график, точки равномерно"""
    c_soft = SOFT_D if dark else SOFT
    c_ink = PAPER_ON_D if dark else INK
    n = max(len(s[2]) for s in series)
    if ymax:
        mx = ymax
    else:
        raw = max(max(s[2]) for s in series) * 1.06
        import math
        mag = 10 ** math.floor(math.log10(raw)) if raw > 0 else 1
        for m in (1.0, 1.25, 1.5, 2.0, 2.5, 3.0, 4.0, 5.0, 6.0, 8.0, 10.0):
            if m * mag >= raw:
                mx = m * mag
                break
        else:
            mx = 10 * mag
        if mx % 4:
            mx = math.ceil(mx / 4.0) * 4.0
    # сетка
    for k in range(5):
        gy = y + h * k / 4.0
        hline(slide, x, gy, w, HAIR_D if dark else HAIR)
        v = mx * (1 - k / 4.0)
        lbl = ylab_fmt(v) if ylab_fmt else f"{v:.0f}"
        tw = text_w(lbl, MONO_MD, 8.5, 0.05)
        block(slide, [lbl], x - tw - 0.16, gy - 0.055, tw + 0.1, MONO_MD, 8.5, c_soft,
              leading=10.5, tracking=0.05)
    step = w / (n - 1) if n > 1 else w
    for name, cl, vals in series:
        pts = [(x + i * step, y + h - (v / mx) * h) for i, v in enumerate(vals)]
        for i in range(len(pts) - 1):
            _seg(slide, pts[i], pts[i + 1], cl, 2.0)
        for px, py in pts:
            dot(slide, px - 0.045, py - 0.045, 0.09, cl)
    if xlabels:
        for i, lb in enumerate(xlabels):
            px = x + i * step
            tw = text_w(lb, MONO_MD, 8.5, 0.06)
            block(slide, [lb], px - tw / 2, y + h + 0.14, tw + 0.12, MONO_MD, 8.5, c_soft,
                  leading=10.5, tracking=0.06)
    # легенда
    ly = y + h + 0.48
    lx = x
    for name, cl, vals in series:
        rect(slide, lx, ly + 0.05, 0.24, 0.045, fill=cl)
        tw = text_w(name, TEXT, size) + 0.02
        block(slide, [name], lx + 0.34, ly - 0.015, tw + 0.1, TEXT, size, c_ink, leading=size * 1.3)
        lx += 0.34 + tw + 0.40
    return ly + 0.36


def _seg(slide, p1, p2, color, width_pt):
    from pptx.util import Emu as _E
    con = slide.shapes.add_connector(1, _E(int(p1[0] * EMU_IN)), _E(int(p1[1] * EMU_IN)),
                                     _E(int(p2[0] * EMU_IN)), _E(int(p2[1] * EMU_IN)))
    con.line.color.rgb = color
    con.line.width = Pt(width_pt)
    return con


def arrow(slide, x, y, w, h, color, direction="right"):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(int(x * EMU_IN)), Emu(int(y * EMU_IN)),
                                 Emu(int(w * EMU_IN)), Emu(int(h * EMU_IN)))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    _noline(shp); _strip_style(shp)
    shp.text_frame.word_wrap = False
    return shp


def ticks(slide, x, y, w, n=24, color=None, dark=False, h=0.055):
    """калибровочная шкала — фирменный мотив «мера»"""
    c = color if color else (HAIR_D if dark else HAIR)
    step = w / n
    for i in range(n + 1):
        hh = h * (1.9 if i % 6 == 0 else 1.0)
        rect(slide, x + i * step, y, 0.008, hh, fill=c)


# ─────────────────────────────────────── фирменный мотив: сигнал измерения
def signal_field(slide, x, y, w, h, n=150, seed=7, base=None, hi=None,
                 hi_from=0.62, hi_to=0.78, bar=0.020):
    """
    Абстрактный сигнал концентрации: поле вертикальных штрихов переменной высоты.
    Читается как график измерения — ровно то, что мы продаём.
    """
    import math
    base = base if base is not None else HAIR_D
    hi = hi if hi is not None else SIGNAL
    step = w / n
    st = seed
    for i in range(n):
        t = i / (n - 1)
        # детерминированный «шум» из трёх гармоник + псевдослучайная добавка
        st = (st * 1103515245 + 12345) % 2147483648
        r = (st / 2147483648.0)
        env = (0.30 + 0.70 * (0.5 + 0.5 * math.sin(t * math.pi * 1.7 - 0.6)))
        v = env * (0.34 + 0.28 * math.sin(t * 21.0) + 0.18 * math.sin(t * 7.3 + 1.1) + 0.42 * r)
        v = max(0.045, min(1.0, v))
        hh = v * h
        inside = hi_from <= t <= hi_to
        rect(slide, x + i * step, y + h - hh, bar, hh, fill=hi if inside else base)
    return y + h


def title_slide(slide, kicker, line1, line2, sub, metrics, footer):
    """титульный слайд"""
    rect(slide, 0, 0, W, H, fill=INK)
    signal_field(slide, ML, 1.02, CW, 1.12, n=132, seed=11,
                 base=RGBColor(0x2B, 0x31, 0x39), hi=SIGNAL, hi_from=0.58, hi_to=0.73)
    hline(slide, ML, 2.14, CW, RGBColor(0x2B, 0x31, 0x39))
    block(slide, [kicker], ML, 2.34, CW, MONO_SB, 9, SIGNAL_L, leading=11.5, tracking=0.20, caps=True)
    y = 2.86
    for ln in (line1, line2):
        sz, lines = fit(ln, DISPLAY_XB, CW * 0.92, 1, 56, 34, 0.5, -0.028)
        y = block(slide, lines, ML, y, CW, DISPLAY_XB, sz, PAPER_ON_D, leading=sz * 1.02, tracking=-0.028)
    y += 0.30
    y = block(slide, wrap(sub, TEXT, 14.5, CW * 0.66), ML, y, CW * 0.7, TEXT, 14.5,
              RGBColor(0xB6, 0xBC, 0xC4), leading=21)
    my = H - MB - 1.02
    hline(slide, ML, my - 0.30, CW, RGBColor(0x2B, 0x31, 0x39))
    kpi_row(slide, metrics, my, dark=True, size=30, label_size=8, gap=0.34)
    block(slide, [footer], ML, H - MB - 0.10, CW, MONO_MD, 8.5, RGBColor(0x6B, 0x72, 0x7B),
          leading=11, tracking=0.10, caps=True)


def section_slide(slide, num, kicker, headline, points, deck=None):
    """разделитель раздела — тёмный, крупный номер"""
    rect(slide, 0, 0, W, H, fill=INK)
    hline(slide, ML, MT, CW, RGBColor(0x2B, 0x31, 0x39))
    block(slide, [kicker], ML, MT + 0.115, CW * 0.6, MONO_SB, 8.5, SIGNAL_L,
          leading=10.5, tracking=0.16, caps=True)
    nx, nw = col(0, 3)
    block(slide, [num], nx, 2.28, nw, DISPLAY_XB, 128, RGBColor(0x24, 0x2A, 0x31),
          leading=118, tracking=-0.04)
    tx, tw_ = col(3, 8)
    sz, lines = fit(headline, DISPLAY_XB, tw_, 3, 40, 26, 0.5, -0.024)
    yy = block(slide, lines, tx, 2.52, tw_, DISPLAY_XB, sz, PAPER_ON_D, leading=sz * 1.10, tracking=-0.024)
    if points:
        yy += 0.42
        hline(slide, tx, yy - 0.20, tw_, RGBColor(0x2B, 0x31, 0x39))
        bullets(slide, points, tx, yy, tw_, dark=True, size=11.5, gap=0.16)


def send_to_back(shape):
    """
    Опускает фигуру в самый низ z-порядка, но НЕ ниже фоновой подложки слайда:
    первая фигура в дереве — это фон, за него уходить нельзя, иначе плашка исчезнет.
    """
    sp = shape._element
    tree = sp.getparent()
    tree.remove(sp)
    first = None
    for ch in tree:
        if ch.tag.endswith('}sp') or ch.tag.endswith('}pic') or ch.tag.endswith('}graphicFrame'):
            first = ch
            break
    if first is None:
        tree.append(sp)
    else:
        first.addnext(sp)
    return shape


def panel_behind(slide, x, y, w, h=None, dark=False, fill=None, border=True, accent=None,
                 accent_w=0.035, bottom=None, pad=0.28, min_h=0.0):
    """
    Плашка, подложенная ПОД уже уложенный текст.

    Предпочтительный способ вызова — bottom=<нижняя граница уложенного текста>:
    высота считается по факту, и текст физически не может вылезти за подложку.
    Передавать h числом можно, но тогда за соответствие отвечает вызывающий —
    именно так в v2 разъехались слайды 11 и 13.
    """
    if bottom is not None:
        h = max(bottom + pad - y, min_h)
    elif h is None:
        raise ValueError("panel_behind: нужен либо bottom=, либо h=")
    shp = panel(slide, x, y, w, h, dark=dark, fill=fill, border=border, accent=None)
    send_to_back(shp)
    if accent is not None:
        a = rect(slide, x, y, accent_w, h, fill=accent)
        send_to_back(a)
        send_to_back(shp)
    return y + h          # нижняя граница подложки — от неё считается всё, что ниже


def card_h(w, tag, head, items, head_size=15, item_size=10.5, pad=0.28, marker="—"):
    """сколько по высоте реально займёт карточка с таким содержимым"""
    iw = w - pad * 2
    yy = pad
    if tag:
        yy += 10.5 / 72.0 + 0.13
    if head:
        hs, hlines = fit(head, DISPLAY_B, iw, 2, head_size, 11, 0.5, -0.015)
        yy += len(hlines) * hs * 1.14 / 72.0 + 0.17
    if items:
        ind = text_w(marker, MONO_MD, item_size, 0) + 0.14
        for it in items:
            h_, rest = (it, None) if isinstance(it, str) else it
            if rest is None:
                n = len(wrap(h_, TEXT, item_size, iw - ind))
            else:
                n = len(mixed_wrap([(h_, TEXT_SB, INK), (rest, TEXT, INK)], iw - ind, item_size))
            yy += n * item_size * 1.40 / 72.0 + 0.11
        yy -= 0.11
    return yy + pad


def row_h(w, specs, **kw):
    """единая высота для ряда карточек — по самой высокой"""
    return max(card_h(w, t, h, i, **kw) for t, h, i in specs)


def qr(slide, data, x, y, size):
    """
    QR-код прямо на слайде.

    Жюри не будет набирать адрес руками. Код рисуется в память и вставляется
    картинкой: отдельного файла в репозитории не появляется, а пересборка
    всегда даёт код, ведущий туда же, куда и подпись рядом.
    """
    import io
    import qrcode
    from qrcode.image.pil import PilImage

    code = qrcode.QRCode(box_size=10, border=1,
                         error_correction=qrcode.constants.ERROR_CORRECT_M)
    code.add_data(data)
    code.make(fit=True)
    # Всегда тёмные модули на белом, даже на тёмном слайде. Инвертированный
    # код красивее вписывается, но часть камер его не берёт, а на защите
    # важнее, чтобы сработало с первого раза.
    img = code.make_image(image_factory=PilImage, fill_color="#1B222B", back_color="#FFFFFF")

    buf = io.BytesIO()
    img.get_image().save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, Emu(int(x * EMU_IN)), Emu(int(y * EMU_IN)),
                             Emu(int(size * EMU_IN)), Emu(int(size * EMU_IN)))
    return y + size
